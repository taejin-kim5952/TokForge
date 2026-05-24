"""RAG 업로드용 문서 → 평문 텍스트 추출."""

from __future__ import annotations

import io
import logging

logger = logging.getLogger(__name__)

RAG_ALLOWED_EXTENSIONS = {
    ".txt",
    ".md",
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".xlsm",
}
RAG_MAX_FILE_BYTES = 10 * 1024 * 1024  # 10MB

# 레거시 바이너리 형식 (.doc, .ppt, .xls) — 별도 변환 필요
RAG_UNSUPPORTED_LEGACY = {".doc", ".ppt", ".xls"}


def file_extension(filename: str | None) -> str:
    if not filename or "." not in filename:
        return ""
    return "." + filename.rsplit(".", 1)[-1].lower()


def allowed_extensions_label() -> str:
    return ", ".join(sorted(RAG_ALLOWED_EXTENSIONS))


def extract_text_from_document(filename: str | None, raw: bytes) -> str:
    ext = file_extension(filename)
    if ext in RAG_UNSUPPORTED_LEGACY:
        raise ValueError(
            f"레거시 형식({ext})은 지원하지 않습니다. "
            f"{ext}x 등 최신 Office 형식으로 저장 후 업로드해 주세요."
        )
    if ext not in RAG_ALLOWED_EXTENSIONS:
        raise ValueError(
            f"지원 형식: {allowed_extensions_label()} (받은 파일: {filename})"
        )

    if len(raw) > RAG_MAX_FILE_BYTES:
        raise ValueError(
            f"파일 크기는 {RAG_MAX_FILE_BYTES // (1024 * 1024)}MB 이하여야 합니다"
        )

    if not raw.strip():
        raise ValueError("빈 파일입니다")

    if ext in (".txt", ".md"):
        text = _extract_plain(raw)
    elif ext == ".pdf":
        text = _extract_pdf(raw)
    elif ext == ".docx":
        text = _extract_docx(raw)
    elif ext == ".pptx":
        text = _extract_pptx(raw)
    elif ext in (".xlsx", ".xlsm"):
        text = _extract_xlsx(raw)
    else:
        raise ValueError(f"지원하지 않는 형식: {ext}")

    text = _normalize(text)
    if not text:
        raise ValueError("문서에서 텍스트를 추출할 수 없습니다 (빈 내용 또는 이미지만 포함)")
    return text


def _normalize(text: str) -> str:
    lines = [ln.rstrip() for ln in text.replace("\r\n", "\n").replace("\r", "\n").split("\n")]
    return "\n".join(lines).strip()


def _extract_plain(raw: bytes) -> str:
    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError:
        raise ValueError("텍스트 파일은 UTF-8 인코딩만 지원합니다")


def _extract_pdf(raw: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as e:
        raise ValueError("PDF 처리 라이브러리가 설치되지 않았습니다 (pypdf)") from e

    reader = PdfReader(io.BytesIO(raw))
    parts: list[str] = []
    for i, page in enumerate(reader.pages):
        chunk = page.extract_text() or ""
        if chunk.strip():
            parts.append(f"--- page {i + 1} ---\n{chunk}")
    return "\n\n".join(parts)


def _extract_docx(raw: bytes) -> str:
    try:
        from docx import Document
    except ImportError as e:
        raise ValueError("Word 처리 라이브러리가 설치되지 않았습니다 (python-docx)") from e

    doc = Document(io.BytesIO(raw))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(raw: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ValueError("PowerPoint 처리 라이브러리가 설치되지 않았습니다 (python-pptx)") from e

    prs = Presentation(io.BytesIO(raw))
    parts: list[str] = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                slide_parts.append(shape.text.strip())
            table = getattr(shape, "table", None)
            if table is not None:
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        slide_parts.append(" | ".join(cells))
        if slide_parts:
            parts.append(f"--- slide {slide_no} ---\n" + "\n".join(slide_parts))
    return "\n\n".join(parts)


def _extract_xlsx(raw: bytes) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as e:
        raise ValueError("Excel 처리 라이브러리가 설치되지 않았습니다 (openpyxl)") from e

    wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_lines: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [
                str(c).strip()
                for c in row
                if c is not None and str(c).strip() != ""
            ]
            if cells:
                sheet_lines.append(" | ".join(cells))
        if sheet_lines:
            parts.append(f"--- sheet: {sheet_name} ---\n" + "\n".join(sheet_lines))
    wb.close()
    return "\n\n".join(parts)


def guess_media_type(filename: str) -> str:
    ext = file_extension(filename)
    return {
        ".txt": "text/plain; charset=utf-8",
        ".md": "text/markdown; charset=utf-8",
        ".pdf": "application/pdf",
        ".docx": (
            "application/vnd.openxmlformats-officedocument"
            ".wordprocessingml.document"
        ),
        ".pptx": (
            "application/vnd.openxmlformats-officedocument"
            ".presentationml.presentation"
        ),
        ".xlsx": (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        ),
        ".xlsm": (
            "application/vnd.ms-excel.sheet.macroEnabled.12"
        ),
    }.get(ext, "application/octet-stream")
