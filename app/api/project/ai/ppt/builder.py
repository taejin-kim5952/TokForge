"""데이터 주도 PPT 빌더 — 템플릿(TemplateContent JSON)을 해석해 .pptx 를 렌더한다.

기존 하드코딩 3종 레이아웃 대신, 템플릿의 blocks[*].elements 를 인터프리터가
요소별로 그린다. build_pptx(deck) 처럼 template 없이 호출하면 DEFAULT_TEMPLATE
(editorial = 기존 디자인)로 폴백하므로 기존 호출부와 호환된다.

바인딩/트랜스폼 계약은 ppt_presets.py 상단 주석 참고.
"""
import io
import logging

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

from .ppt_presets import DEFAULT_TEMPLATE

logger = logging.getLogger(__name__)

FONT_BODY = "Pretendard"        # 폰트 ref 해석 실패 시 최종 폴백
_FALLBACK_COLOR = "9AA0AD"
_BLANK = 6


# ─────────────────────────────────────────────
#  Low-level helpers (기존 유지)
# ─────────────────────────────────────────────
def _bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def _rect(slide, left, top, width, height, color, line=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if not line:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _line(slide, left, top, width, color, thickness_pt=0.75):
    """얇은 가로 룰선 — 두께를 Pt로 정확히 컨트롤."""
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Emu(int(Pt(thickness_pt)))
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _box(slide, left, top, width, height, anchor=None):
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    return tf


def _para(tf, text, *, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT,
          space_after=6, line_spacing=None, font=None, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    if line_spacing is not None:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font or FONT_BODY
    return p


def _run(paragraph, text, *, size, color, bold=False, italic=False, font=None):
    r = paragraph.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font or FONT_BODY
    return r


# ─────────────────────────────────────────────
#  Resolve helpers (theme ref / enum / binding)
# ─────────────────────────────────────────────
_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
_ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
_ROMANS = ["", "I", "II", "III", "IV", "V", "VI", "VII", "VIII", "IX", "X",
           "XI", "XII", "XIII", "XIV", "XV"]


def _roman(n):
    return _ROMANS[n] + "." if 0 < n < len(_ROMANS) else f"{n:02d}."


def _color(theme, ref):
    hexv = (theme.get("colors", {}) or {}).get(ref)
    if not hexv:
        logger.warning("unknown color ref %r", ref)
        hexv = _FALLBACK_COLOR
    return RGBColor.from_string(hexv.lstrip("#"))


def _font(theme, ref):
    return (theme.get("fonts", {}) or {}).get(ref) or FONT_BODY


def _align_of(v):
    return _ALIGN.get(v, PP_ALIGN.LEFT)


def _anchor_of(v):
    return _ANCHOR.get(v)


def _apply_transforms(s, transforms):
    for tr in transforms or []:
        if tr == "upper":
            s = s.upper()
        elif tr == "lower":
            s = s.lower()
        elif tr == "trim":
            s = s.strip()
        elif tr == "pad2":
            try:
                s = f"{int(s):02d}"
            except (TypeError, ValueError):
                pass
        elif tr == "roman":
            try:
                s = _roman(int(s))
            except (TypeError, ValueError):
                pass
    return s


def _resolve_bind(bind, slide, ctx, item=None, index=None):
    if not bind:
        return ""
    if bind.startswith("deck."):
        return str(ctx["deck"].get(bind[5:], "") or "")
    if bind == "@pageNo":
        return str(ctx["page_no"])
    if bind == "@total":
        return str(ctx["total"])
    if bind == "@sectionNo":
        return str(ctx["section_no"])
    if bind == "@sectionTitle":
        return str(ctx.get("section_title") or "")
    if bind == "@index":
        return str(0 if index is None else index)
    if bind == "@itemNo":
        return str(1 if index is None else index + 1)
    if bind == "@item":
        if item is None or isinstance(item, dict):
            return ""
        return str(item)
    if bind.startswith("@item."):
        return str(item.get(bind[6:], "")) if isinstance(item, dict) else ""
    val = slide.get(bind, "")
    return "" if val is None else str(val)


# ─────────────────────────────────────────────
#  Element interpreter
# ─────────────────────────────────────────────
def _style_para(tf, text, theme, style, *, first=True):
    _para(
        tf, text,
        size=style.get("size", 18),
        color=_color(theme, style.get("color", "ink")),
        bold=style.get("bold", False),
        italic=style.get("italic", False),
        align=_align_of(style.get("align")),
        space_after=style.get("spaceAfter", 6),
        line_spacing=style.get("lineSpacing"),
        font=_font(theme, style.get("font", "body")),
        first=first,
    )


def _render_text(slide, el, kind, theme, slide_data, ctx):
    g = el["geom"]
    style = el.get("style", {})
    if kind == "text":
        text = _apply_transforms(el.get("text", ""), el.get("transform"))
    elif kind == "bound":
        text = _apply_transforms(
            _resolve_bind(el["bind"], slide_data, ctx), el.get("transform")
        )
        if not text:
            text = el.get("fallback", "")
    else:  # pageNumber
        text = (el.get("format", "")
                .replace("{page:02d}", f"{ctx['page_no']:02d}")
                .replace("{total:02d}", f"{ctx['total']:02d}"))
    if not text:
        return
    tf = _box(slide, Inches(g["x"]), Inches(g["y"]), Inches(g["w"]), Inches(g["h"]),
              anchor=_anchor_of(style.get("anchor")))
    _style_para(tf, text, theme, style)


def _render_repeater(slide, el, theme, slide_data, ctx):
    arr = slide_data.get(el.get("bind", ""), [])
    if not isinstance(arr, list):
        return
    maxn = el.get("maxItems")
    if maxn:
        arr = arr[:maxn]
    item = el.get("item", {})
    if el.get("render", "paragraphs") == "paragraphs":
        _render_repeater_paragraphs(slide, el, arr, item, theme, slide_data, ctx)
    else:
        _render_repeater_boxes(slide, el, arr, item, theme, slide_data, ctx)


def _render_repeater_paragraphs(slide, el, arr, item, theme, slide_data, ctx):
    g = el["geom"]
    tf = _box(slide, Inches(g["x"]), Inches(g["y"]), Inches(g["w"]), Inches(g["h"]))
    ps = item.get("paraStyle", {})
    runs = [e for e in item.get("elements", []) if e.get("kind") == "run"]
    for i, it in enumerate(arr):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = _align_of(ps.get("align"))
        p.space_after = Pt(ps.get("spaceAfter", 6))
        if ps.get("lineSpacing") is not None:
            p.line_spacing = ps["lineSpacing"]
        for run in runs:
            st = run.get("style", {})
            if run.get("bind"):
                txt = _resolve_bind(run["bind"], slide_data, ctx, item=it, index=i)
            elif run.get("text") is not None:
                txt = run["text"]
            else:
                txt = ""
            txt = _apply_transforms(txt, run.get("transform"))
            _run(p, txt,
                 size=st.get("size", 18),
                 color=_color(theme, st.get("color", "ink")),
                 bold=st.get("bold", False),
                 italic=st.get("italic", False),
                 font=_font(theme, st.get("font", "body")))


def _render_repeater_boxes(slide, el, arr, item, theme, slide_data, ctx):
    g = el["geom"]
    item_geom = item.get("geom", {"x": 0, "y": 0, "w": g["w"], "h": 1})
    direction = el.get("direction", "vertical")
    columns = max(1, int(el.get("columns", 1)))
    stride = el.get("stride", item_geom["h"])
    spacing = el.get("spacing", 0.0)
    overflow = el.get("overflow")
    cell_w = g["w"] / columns if direction == "grid" else item_geom["w"]

    for i, it in enumerate(arr):
        if direction == "grid":
            ox = g["x"] + (i % columns) * cell_w
            oy = g["y"] + (i // columns) * (stride + spacing)
        elif direction == "horizontal":
            ox = g["x"] + i * (item_geom["w"] + spacing)
            oy = g["y"]
        else:  # vertical
            ox = g["x"]
            oy = g["y"] + i * (stride + spacing)
        if overflow == "clip" and oy + item_geom["h"] > g["y"] + g["h"] + 0.01:
            break
        for sub in item.get("elements", []):
            _render_box_child(slide, sub, ox, oy, theme, slide_data, ctx, it, i)


def _render_box_child(slide, sub, ox, oy, theme, slide_data, ctx, item, index):
    kind = sub.get("kind")
    g = sub.get("geom", {"x": 0, "y": 0, "w": 1, "h": 1})
    ax = ox + g.get("x", 0)
    ay = oy + g.get("y", 0)
    if kind == "rect":
        _rect(slide, Inches(ax), Inches(ay), Inches(g["w"]), Inches(g["h"]),
              _color(theme, sub["fill"]))
        return
    if kind == "line":
        _line(slide, Inches(ax), Inches(ay), Inches(g["w"]),
              _color(theme, sub["color"]), sub.get("thicknessPt", 0.75))
        return
    if kind in ("text", "bound", "run"):
        style = sub.get("style", {})
        if kind == "text":
            txt = _apply_transforms(sub.get("text", ""), sub.get("transform"))
        else:
            txt = _apply_transforms(
                _resolve_bind(sub.get("bind", ""), slide_data, ctx, item, index),
                sub.get("transform"),
            )
            if not txt:
                txt = sub.get("fallback", "")
        if not txt:
            return
        tf = _box(slide, Inches(ax), Inches(ay), Inches(g["w"]), Inches(g["h"]),
                  anchor=_anchor_of(style.get("anchor")))
        _style_para(tf, txt, theme, style)


def _render_image(slide, el, ctx):
    """이미지 에셋 렌더 — (project_id, asset)로 디스크 경로 해석 후 add_picture."""
    asset_id = el.get("asset")
    pid = ctx.get("project_id")
    if not asset_id or pid is None:
        return
    # 지연 import — builder↔service 순환 및 __main__ 단독 실행 보호
    from app.services.ppt_template_asset_service import resolve_asset_path

    path = resolve_asset_path(pid, asset_id)
    if not path or not path.is_file():
        logger.warning("image asset missing: project=%s asset=%s", pid, asset_id)
        return
    g = el["geom"]
    slide.shapes.add_picture(
        str(path), Inches(g["x"]), Inches(g["y"]), Inches(g["w"]), Inches(g["h"])
    )


def _render_element(slide, el, theme, slide_data, ctx):
    vif = el.get("visibleIf")
    if vif and not slide_data.get(vif):
        return
    kind = el.get("kind")
    if kind == "rect":
        g = el["geom"]
        _rect(slide, Inches(g["x"]), Inches(g["y"]), Inches(g["w"]), Inches(g["h"]),
              _color(theme, el["fill"]))
    elif kind == "line":
        g = el["geom"]
        _line(slide, Inches(g["x"]), Inches(g["y"]), Inches(g["w"]),
              _color(theme, el["color"]), el.get("thicknessPt", 0.75))
    elif kind in ("text", "bound", "pageNumber"):
        _render_text(slide, el, kind, theme, slide_data, ctx)
    elif kind == "repeater":
        _render_repeater(slide, el, theme, slide_data, ctx)
    elif kind == "image":
        _render_image(slide, el, ctx)
    else:
        logger.warning("unknown element kind %r", kind)


# ─────────────────────────────────────────────
#  Builder
# ─────────────────────────────────────────────
def build_pptx(deck, template: dict | None = None, *, project_id: int | None = None) -> bytes:
    tpl = template or DEFAULT_TEMPLATE
    prs = Presentation()
    sl = tpl.get("slide", {})
    prs.slide_width = Inches(sl.get("w", 13.333))
    prs.slide_height = Inches(sl.get("h", 7.5))

    theme = tpl.get("theme", {})
    blocks = tpl.get("blocks", [])
    by_type = {b["type"]: b for b in blocks}
    fallback_block = by_type.get("bullets") or (blocks[0] if blocks else None)

    slides = deck["slides"]
    total = len(slides)
    section_no = 0
    section_title = None

    for i, sd in enumerate(slides):
        stype = sd.get("type")
        if stype == "section":
            section_no += 1
            section_title = sd.get("title", "")
        block = by_type.get(stype) or fallback_block
        if not block:
            continue
        ctx = {
            "deck": deck,
            "page_no": i + 1,
            "total": total,
            "section_no": section_no,
            "section_title": section_title,
            "project_id": project_id,
        }
        s = prs.slides.add_slide(prs.slide_layouts[_BLANK])
        _bg(s, _color(theme, block.get("bg", "paper")))
        for el in block.get("elements", []):
            try:
                _render_element(s, el, theme, sd, ctx)
            except Exception:
                logger.warning("element render failed (kind=%r)", el.get("kind"),
                               exc_info=True)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ─────────────────────────────────────────────
#  실행  (python builder.py → 샘플 .pptx 생성)
# ─────────────────────────────────────────────
if __name__ == "__main__":
    DECK = {
        "title": "TokForge 소개",
        "slides": [
            {"type": "title", "title": "TokForge",
             "subtitle": "셀프호스팅 LLM 토큰 절감 워크스페이스", "meta": "2026 PROPOSAL"},
            {"type": "section", "title": "배경 및 필요성", "note": "기존 LLM 운영의 한계 진단"},
            {"type": "bullets", "title": "기존 LLM 운영의 과제", "bullets": [
                "유료 API 토큰 비용이 사용량에 비례해 선형 증가",
                "캐시·로컬 모델 활용 부재로 동일 요청 반복 과금",
                "어떤 호출이 비용을 유발하는지 가시성이 없음",
                "프로젝트별 컨텍스트 격리·재사용 체계 미비",
            ]},
            {"type": "bullets", "title": "TokForge의 접근",
             "subtitle": "캐시·로컬 우선 라우팅 + 비용 가시화",
             "bullets": [
                "캐시 히트 + 로컬 Ollama 우선 라우팅으로 유료 호출 최소화",
                "요청별 토큰·비용 지표 기록으로 절감액 산출",
                "프로젝트 단위 컨텍스트 격리로 재사용성 확보",
            ]},
            {"type": "section", "title": "기대 효과", "note": "비용·운영 관점의 가치"},
            {"type": "bullets", "title": "비용 절감 ROI", "bullets": [
                "캐시 히트·로컬 처리 시 해당 호출 비용 100% 절감",
                "유료 호출은 실제 발생 비용만 베이스라인 대비 집계",
            ]},
        ],
    }
    data = build_pptx(DECK)
    out = r"d:\tuning\_design_sample.pptx"
    with open(out, "wb") as f:
        f.write(data)
    print(f"OK {len(data)} bytes, {len(DECK['slides'])} slides -> {out}")
