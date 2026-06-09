"""업로드된 회사 .pptx → TemplateContent best-effort 변환 파서.

순수 함수: 바이트 입력 → {content, assets, report} (디스크/DB 미접촉).
- 실제 슬라이드를 순회해 블록(title/section/bullets)으로 분류
- 도형 → 요소(bound/text/repeater/image/rect/line) 매핑, 기하는 EMU→inch
- 테마(색/폰트)는 관측된 값으로 합성 (python-pptx 에 테마 공개 API 없음)
- 이미지는 hint("__ASSET_i__")로 표시 → 엔드포인트가 저장 후 실제 asset id 로 치환

분류/매핑 휴리스틱은 실제 회사 샘플로 튜닝하는 표면이다 (report 참고).
"""
import io
import logging
from collections import Counter

from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

logger = logging.getLogger(__name__)

_EMU_PER_IN = 914400


def parse_pptx_to_template(raw: bytes) -> dict:
    prs = Presentation(io.BytesIO(raw))
    report = {"slides_seen": 0, "blocks": 0, "skipped": [], "warnings": [],
              "fonts": 0, "colors": 0}
    assets: list[dict] = []

    # 슬라이드 치수
    sw = round((prs.slide_width or 12192000) / _EMU_PER_IN, 3)
    sh = round((prs.slide_height or 6858000) / _EMU_PER_IN, 3)
    margin = round(min(sw, sh) * 0.06, 2)

    # ── 패스1: 색/폰트 빈도 → 시드 결정 ──
    font_ct: Counter = Counter()
    color_ct: Counter = Counter()
    for slide in prs.slides:
        for shp in slide.shapes:
            try:
                if not shp.has_text_frame:
                    continue
                for p in shp.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.name:
                            font_ct[r.font.name] += 1
                        hx = _run_hex(r)
                        if hx:
                            color_ct[hx] += 1
            except Exception:
                continue
    body_font = font_ct.most_common(1)[0][0] if font_ct else "Pretendard"
    ink_hex = color_ct.most_common(1)[0][0] if color_ct else "#0F172A"

    # ref 맵 (hex/name → ref), 시드 ink/paper/body
    colors: dict[str, str] = {}
    fonts: dict[str, str] = {}
    colors[_norm_hex(ink_hex)] = "ink"
    colors.setdefault("#FFFFFF", "paper")
    fonts[body_font] = "body"
    ctr = {"c": 1, "f": 1}

    def reg_color(hx: str) -> str:
        key = _norm_hex(hx)
        if key in colors:
            return colors[key]
        ref = f"c{ctr['c']}"
        ctr["c"] += 1
        colors[key] = ref
        return ref

    def reg_font(name: str) -> str:
        if name in fonts:
            return fonts[name]
        ref = f"f{ctr['f']}"
        ctr["f"] += 1
        fonts[name] = ref
        return ref

    # ── 패스2: 슬라이드 → 블록 ──
    blocks: list[dict] = []
    seen_sigs: set = set()
    used_types: dict[str, int] = {}

    for slide in prs.slides:
        report["slides_seen"] += 1
        layout = slide.slide_layout
        elements: list[dict] = []
        for shp in slide.shapes:
            try:
                els = _shape_to_elements(shp, layout, reg_color, reg_font, assets, report)
                elements.extend(els)
            except Exception as e:  # noqa: BLE001
                report["warnings"].append(
                    f"shape '{getattr(shp, 'name', '?')}' 처리 실패: {e}"
                )
        if not elements:
            report["skipped"].append(f"slide {report['slides_seen']}: 빈 슬라이드/미지원")
            continue

        sig = tuple((e["kind"], round(e["geom"]["x"], 1), round(e["geom"]["y"], 1))
                    for e in elements)
        if sig in seen_sigs:
            report["skipped"].append(f"slide {report['slides_seen']}: 유사 레이아웃 중복")
            continue
        seen_sigs.add(sig)

        btype = _classify(slide)
        if btype in used_types:
            used_types[btype] += 1
            btype = f"{btype}{used_types[btype]}"
        else:
            used_types[btype] = 1
        blocks.append({
            "type": btype,
            "label": _slide_title_text(slide) or btype.title(),
            "bg": "paper",
            "elements": elements,
        })

    report["blocks"] = len(blocks)
    report["fonts"] = len(fonts)
    report["colors"] = len(colors)

    content = {
        "schema_version": 1,
        "slide": {"w": sw, "h": sh, "margin": margin},
        "theme": {
            "colors": {ref: hx for hx, ref in colors.items()},
            "fonts": {ref: name for name, ref in fonts.items()},
        },
        "blocks": blocks,
    }
    return {"content": content, "assets": assets, "report": report}


# ───────── 분류/추출 헬퍼 ─────────
def _shape_to_elements(shape, layout, reg_color, reg_font, assets, report) -> list[dict]:
    # 1) 그림
    if _is_picture(shape):
        try:
            img = shape.image
            hint = f"__ASSET_{len(assets)}__"
            assets.append({"hint": hint, "ext": img.ext, "blob": img.blob})
            return [{"kind": "image", "geom": _geom(shape, layout, report),
                     "asset": hint, "fit": "contain"}]
        except Exception as e:  # noqa: BLE001
            report["skipped"].append(f"picture 추출 실패: {e}")
            return []

    # 2) 텍스트 (placeholder/텍스트박스/오토셰이프 텍스트)
    if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
        geom = _geom(shape, layout, report)
        style = _style_from(shape, reg_color, reg_font)
        pt = _placeholder_type(shape)
        if pt in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            return [_bound("title", geom, style)]
        if pt == PP_PLACEHOLDER.SUBTITLE:
            return [_bound("subtitle", geom, style, visible_if="subtitle")]
        paras = [p for p in shape.text_frame.paragraphs if p.text.strip()]
        if pt in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT) or len(paras) >= 2:
            return [{
                "kind": "repeater", "bind": "bullets", "geom": geom,
                "direction": "vertical", "stride": 0.0, "render": "paragraphs",
                "item": {
                    "paraStyle": {"spaceAfter": 10,
                                  "lineSpacing": style.get("lineSpacing", 1.4),
                                  "align": style.get("align", "left")},
                    "elements": [{"kind": "run", "bind": "@item", "style": style}],
                },
            }]
        # 단문 정적 텍스트 (로고 텍스트/태그라인 등) → 리터럴 보존
        return [{"kind": "text", "geom": geom,
                 "text": shape.text_frame.text.strip(), "style": style}]

    # 3) 단색 채움 오토셰이프 → rect / line
    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        fill = _fill_hex(shape)
        if fill:
            geom = _geom(shape, layout, report)
            if geom["h"] < 0.06:
                return [{"kind": "line", "geom": geom, "color": reg_color(fill),
                         "thicknessPt": max(0.5, round(geom["h"] * 72, 2))}]
            return [{"kind": "rect", "geom": geom, "fill": reg_color(fill)}]

    # 4) 미지원 (그룹/SmartArt/차트/표/커넥터/미디어 등)
    report["skipped"].append(
        f"미지원 도형 '{getattr(shape, 'name', '?')}' (type={shape.shape_type})"
    )
    return []


def _classify(slide) -> str:
    has_title = has_ctr = has_sub = False
    body = 0
    for shp in slide.shapes:
        pt = _placeholder_type(shp)
        if pt == PP_PLACEHOLDER.TITLE:
            has_title = True
        elif pt == PP_PLACEHOLDER.CENTER_TITLE:
            has_ctr = True
        elif pt == PP_PLACEHOLDER.SUBTITLE:
            has_sub = True
        elif pt in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            body += 1
    if has_ctr or has_sub:
        return "title"
    if has_title and body == 0:
        return "section"
    return "bullets"


# ───────── 저수준 헬퍼 ─────────
def _emu_in(v) -> float:
    return round(v / _EMU_PER_IN, 3)


def _geom(shape, layout, report) -> dict:
    l, t, w, h = shape.left, shape.top, shape.width, shape.height
    if None in (l, t, w, h) and getattr(shape, "is_placeholder", False):
        ph = _layout_ph(layout, shape.placeholder_format.idx)
        if ph is not None:
            l = l if l is not None else ph.left
            t = t if t is not None else ph.top
            w = w if w is not None else ph.width
            h = h if h is not None else ph.height
    if None in (l, t, w, h):
        report["warnings"].append(
            f"shape '{getattr(shape, 'name', '?')}' 기하 미상 → 기본값"
        )
        l, t, w, h = (l or 0), (t or 0), (w or 914400), (h or 457200)
    return {"x": _emu_in(l), "y": _emu_in(t), "w": _emu_in(w), "h": _emu_in(h)}


def _layout_ph(layout, idx):
    try:
        for ph in layout.placeholders:
            if ph.placeholder_format.idx == idx:
                return ph
    except Exception:
        pass
    return None


def _placeholder_type(shape):
    try:
        if shape.is_placeholder:
            return shape.placeholder_format.type
    except Exception:
        pass
    return None


def _is_picture(shape) -> bool:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True
    try:
        _ = shape.image
        return True
    except Exception:
        return False


def _first_run(shape):
    try:
        for p in shape.text_frame.paragraphs:
            if p.runs:
                return p, p.runs[0]
        for p in shape.text_frame.paragraphs:
            if p.text.strip():
                return p, None
    except Exception:
        pass
    return None, None


def _style_from(shape, reg_color, reg_font) -> dict:
    p, r = _first_run(shape)
    style = {"font": "body", "size": 18, "color": "ink", "align": "left"}
    if p is not None:
        a = p.alignment
        if a == PP_ALIGN.CENTER:
            style["align"] = "center"
        elif a == PP_ALIGN.RIGHT:
            style["align"] = "right"
    try:
        va = shape.text_frame.vertical_anchor
        if va == MSO_ANCHOR.MIDDLE:
            style["anchor"] = "middle"
        elif va == MSO_ANCHOR.BOTTOM:
            style["anchor"] = "bottom"
    except Exception:
        pass
    if r is not None:
        f = r.font
        if f.name:
            style["font"] = reg_font(f.name)
        try:
            if f.size is not None:
                style["size"] = round(f.size.pt, 1)
        except Exception:
            pass
        if f.bold:
            style["bold"] = True
        if f.italic:
            style["italic"] = True
        hx = _run_hex(r)
        if hx:
            style["color"] = reg_color(hx)
    return style


def _run_hex(run) -> str | None:
    try:
        c = run.font.color
        if c is not None and c.type is not None:
            return "#" + str(c.rgb)  # 테마색 참조면 .rgb 가 예외 → None
    except Exception:
        return None
    return None


def _fill_hex(shape) -> str | None:
    try:
        if shape.fill.type == MSO_FILL.SOLID:
            return "#" + str(shape.fill.fore_color.rgb)
    except Exception:
        return None
    return None


def _norm_hex(hx: str) -> str:
    s = (hx or "").strip().lstrip("#").upper()
    return "#" + s if len(s) == 6 else "#000000"


def _bound(bind, geom, style, visible_if=None) -> dict:
    d = {"kind": "bound", "geom": geom, "bind": bind, "style": style}
    if visible_if is not None:
        d["visibleIf"] = visible_if
    return d


def _slide_title_text(slide) -> str:
    try:
        if slide.shapes.title and slide.shapes.title.text.strip():
            return slide.shapes.title.text.strip()[:40]
    except Exception:
        pass
    return ""
